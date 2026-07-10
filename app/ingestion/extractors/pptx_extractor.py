from pptx import Presentation

from app.ingestion.base_extractor import BaseExtractor
from app.models.knowledge_source import KnowledgeSource

class PPTXExtractor(BaseExtractor):
    def extract(self, source: KnowledgeSource):
        prs = Presentation(source.metadata["path"])

        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): # "Hey, does this specific object even have a 'text' property?"
                    text.append(getattr(shape, "text"))

        source.raw_content = "\n".join(text)
        return source